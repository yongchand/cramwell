from dotenv import load_dotenv
import pandas as pd
import json
import os
import uuid
import warnings
from datetime import datetime
import re
from jinja2 import Template, Environment, FileSystemLoader
from pathlib import Path

from pydantic import BaseModel, Field, model_validator
from openai import OpenAI
from openai.types.chat import ChatCompletionMessageParam as ChatMessage
from .pinecone_service import pinecone_service
from typing_extensions import override
from typing import List, Tuple, Union, Optional, Dict, cast
from typing_extensions import Self
from .database import supabase


load_dotenv()

# Initialize Jinja2 environment for templating
def get_template_environment():
    """Get Jinja2 template environment."""
    template_dir = Path(__file__).parent / "prompts"
    return Environment(loader=FileSystemLoader(template_dir))

def format_response_with_template(raw_response: str, question: str) -> str:
    """
    Format the raw response using the response template.
    """
    try:
        env = get_template_environment()
        template = env.get_template("response_template.jinja")
        
        # Render the template with the response data
        formatted_response = template.render(
            question=question,
            raw_response=raw_response
        )
        
        return formatted_response
        
    except Exception as e:
        # Fallback to simple formatting
        return f"**Answer:**\n\n{raw_response}\n\n---\n\n*This response is based on your uploaded documents.*"








class ClaimVerification(BaseModel):
    claim_is_true: bool = Field(
        description="Based on the provided sources information, the claim passes or not."
    )
    supporting_citations: Optional[List[str]] = Field(
        description="A minimum of one and a maximum of three citations from the sources supporting the claim. If the claim is not supported, please leave empty",
        default=None,
        min_length=1,
        max_length=3,
    )

    @model_validator(mode="after")
    def validate_claim_ver(self) -> Self:
        if not self.claim_is_true and self.supporting_citations is not None:
            self.supporting_citations = ["The claim was deemed false."]
        return self


# Initialize OpenAI client
openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# Initialize LLM for structured output (simplified for now)
LLM_STRUCT = openai_client
LLM_VERIFIER = openai_client

# Cache DocumentConverter at module level with memory-optimized settings
try:
    from docling.document_converter import DocumentConverter
    from docling.datamodel.pipeline_options import PdfPipelineOptions
    from docling.datamodel.base_models import InputFormat
    from docling.backend.pypdfium2_backend import PyPdfiumDocumentBackend
    from docling.document_converter import (DocumentConverter,PdfFormatOption,WordFormatOption)

    pipeline_options = PdfPipelineOptions()
    pipeline_options.do_ocr = False
    pipeline_options.do_table_structure = True
    # Use memory-optimized settings for fallback converter
    converter = DocumentConverter(
        format_options={
            InputFormat.PDF: PdfFormatOption(
                                pipeline_options=pipeline_options, backend=PyPdfiumDocumentBackend), # pipeline options go here.
        }
    )
    DOC_CONVERTER = converter
except Exception as e:
    DOC_CONVERTER = None


async def parse_file_pymupdf(
    file_path: str, with_images: bool = False, with_tables: bool = False
) -> Union[Tuple[Optional[str], Optional[List[str]], Optional[List[pd.DataFrame]]]]:
    """
    Parse a file using PyMuPDF for lightweight document processing.
    """
    images: Optional[List[str]] = None
    text: Optional[str] = None
    tables: Optional[List[pd.DataFrame]] = None
    
    try:
        import fitz  # PyMuPDF
        
        # Open and extract text
        doc = fitz.open(file_path)
        text = ""
        
        for page in doc:
            text += page.get_text()
        
        # Extract images if requested
        if with_images:
            import base64
            images = []
            for page_num, page in enumerate(doc):
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        
                        # Convert to base64
                        img_data = pix.tobytes("png")
                        img_base64 = base64.b64encode(img_data).decode('utf-8')
                        images.append(f"data:image/png;base64,{img_base64}")
                        
                        pix = None  # Free memory
                    except Exception as e:
                        # Skip images that can't be extracted
                        pass
        
        doc.close()
        
        # Extract tables if requested (simplified)
        if with_tables:
            tables = []
        
        return text, images, tables
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return None, None, None
    finally:
        # Force garbage collection
        import gc
        gc.collect()


async def parse_file_docling(
    file_path: str, with_images: bool = False, with_tables: bool = False
) -> Union[Tuple[Optional[str], Optional[List[str]], Optional[List[pd.DataFrame]]]]:
    """
    Parse a file using Docling for advanced document processing (OCR, tables, etc.).
    """
    images: Optional[List[str]] = None
    text: Optional[str] = None
    tables: Optional[List[pd.DataFrame]] = None
    
    try:
        # Use cached DocumentConverter
        if DOC_CONVERTER is None:
            from docling.document_converter import DocumentConverter
            from docling.datamodel.pipeline_options import PdfPipelineOptions
            from docling.datamodel.base_models import InputFormat
            from docling.backend.pypdfium2_backend import PyPdfiumDocumentBackend
            from docling.document_converter import (DocumentConverter,PdfFormatOption,WordFormatOption)

            pipeline_options = PdfPipelineOptions()
            pipeline_options.do_ocr = False
            pipeline_options.do_table_structure = True
            # Use memory-optimized settings for fallback converter
            converter = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(
                                        pipeline_options=pipeline_options, backend=PyPdfiumDocumentBackend), # pipeline options go here.
                }
            )
        else:
            converter = DOC_CONVERTER
        result = converter.convert(file_path)
        # Extract text content from markdown
        text = result.document.export_to_markdown()
        # Extract tables if requested
        if with_tables:
            tables = []
        # Extract images if requested
        if with_images:
            images = []
        return text, images, tables
    except Exception as e:
        import traceback
        traceback.print_exc()
        return None, None, None
    finally:
        # Force garbage collection
        import gc
        gc.collect()


def is_handwritten_or_poor_extraction(text: str) -> bool:
    """
    Check if the extracted text suggests handwritten content or poor extraction.
    """
    if not text or len(text.strip()) < 50:
        return True  # Too little text extracted
    
    # Check for common handwritten indicators
    handwritten_indicators = [
        "handwritten", "handwriting", "scanned", "image", "photo",
        "unreadable", "illegible", "blurry", "fuzzy"
    ]
    
    text_lower = text.lower()
    for indicator in handwritten_indicators:
        if indicator in text_lower:
            return True
    
    # Check if text looks like OCR output (lots of random characters)
    if len(text) > 100:
        # Count non-alphanumeric characters
        non_alphanumeric = sum(1 for c in text if not c.isalnum() and not c.isspace())
        ratio = non_alphanumeric / len(text)
        if ratio > 0.3:  # More than 30% non-alphanumeric suggests OCR issues
            return True
    
    return False


async def parse_spreadsheet_file(
    file_path: str, with_tables: bool = False
) -> Union[Tuple[Optional[str], Optional[List[str]], Optional[List[pd.DataFrame]]]]:
    """
    Parse Excel (.xlsx) and CSV files.
    """
    text: Optional[str] = None
    tables: Optional[List[pd.DataFrame]] = []
    images: Optional[List[str]] = None
    
    try:
        import pandas as pd
        
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == '.xlsx':
            # Read Excel file - get all sheets
            excel_file = pd.ExcelFile(file_path)
            sheet_names = excel_file.sheet_names
            
            text_parts = []
            for sheet_name in sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text_parts.append(f"Sheet: {sheet_name}\n")
                text_parts.append(df.to_string(index=False))
                text_parts.append("\n\n")
                
                if with_tables:
                    tables.append(df)
            
            text = "".join(text_parts)
            
        elif file_ext == '.csv':
            # Read CSV file
            df = pd.read_csv(file_path)
            text = f"CSV Data:\n{df.to_string(index=False)}"
            
            if with_tables:
                tables.append(df)
        
        return text, images, tables
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return None, None, None
    finally:
        # Force garbage collection
        import gc
        gc.collect()


async def parse_jupyter_notebook(
    file_path: str
) -> Union[Tuple[Optional[str], Optional[List[str]], Optional[List[pd.DataFrame]]]]:
    """
    Parse Jupyter notebook (.ipynb) files.
    """
    text: Optional[str] = None
    tables: Optional[List[pd.DataFrame]] = None
    images: Optional[List[str]] = None
    
    try:
        import json
        
        with open(file_path, 'r', encoding='utf-8') as f:
            notebook_data = json.load(f)
        
        text_parts = []
        
        for cell in notebook_data.get('cells', []):
            cell_type = cell.get('cell_type', '')
            source = ''.join(cell.get('source', []))
            
            if cell_type == 'markdown':
                text_parts.append(f"# Markdown Cell\n{source}\n\n")
            elif cell_type == 'code':
                text_parts.append(f"# Code Cell\n```python\n{source}\n```\n\n")
                
                # Add output if available
                outputs = cell.get('outputs', [])
                for output in outputs:
                    if output.get('output_type') == 'execute_result':
                        data = output.get('data', {})
                        if 'text/plain' in data:
                            text_parts.append(f"Output: {data['text/plain']}\n\n")
        
        text = "".join(text_parts)
        return text, images, tables
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return None, None, None
    finally:
        # Force garbage collection
        import gc
        gc.collect()


async def parse_powerpoint_file(
    file_path: str
) -> Union[Tuple[Optional[str], Optional[List[str]], Optional[List[pd.DataFrame]]]]:
    """
    Parse PowerPoint (.ppt, .pptx) files.
    """
    text: Optional[str] = None
    tables: Optional[List[pd.DataFrame]] = None
    images: Optional[List[str]] = None
    
    try:
        from pptx import Presentation
        import base64
        
        prs = Presentation(file_path)
        text_parts = []
        images = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"Slide {slide_num}:\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(f"{shape.text}\n")
                
                # Extract images from shapes
                if hasattr(shape, "image"):
                    try:
                        image = shape.image
                        # Convert image to base64 for storage
                        image_bytes = image.blob
                        image_ext = image.ext
                        image_data = base64.b64encode(image_bytes).decode('utf-8')
                        images.append(f"data:image/{image_ext};base64,{image_data}")
                    except Exception as e:
                        # Skip images that can't be extracted
                        pass
            
            text_parts.append("\n")
        
        text = "".join(text_parts)
        return text, images, tables
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return None, None, None
    finally:
        # Force garbage collection
        import gc
        gc.collect()


async def parse_docx_file(
    file_path: str
) -> Union[Tuple[Optional[str], Optional[List[str]], Optional[List[pd.DataFrame]]]]:
    """
    Parse Word (.docx) files using python-docx.
    """
    text: Optional[str] = None
    tables: Optional[List[pd.DataFrame]] = None
    images: Optional[List[str]] = None
    
    try:
        from docx import Document
        
        doc = Document(file_path)
        text_parts = []
        
        # Extract text from paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text.strip() + "\n")
        
        # Extract text from tables
        for table in doc.tables:
            text_parts.append("\n--- Table ---\n")
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text.strip())
                text_parts.append(" | ".join(row_text) + "\n")
            text_parts.append("--- End Table ---\n\n")
        
        # Extract text from headers and footers
        for section in doc.sections:
            header = section.header
            for paragraph in header.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(f"Header: {paragraph.text.strip()}\n")
            
            footer = section.footer
            for paragraph in footer.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(f"Footer: {paragraph.text.strip()}\n")
        
        text = "".join(text_parts)
        return text, images, tables
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return None, None, None
    finally:
        # Force garbage collection
        import gc
        gc.collect()


async def parse_file(
    file_path: str, with_images: bool = False, with_tables: bool = False
) -> Union[Tuple[Optional[str], Optional[List[str]], Optional[List[pd.DataFrame]]]]:
    """
    Smart hybrid parsing: Use specialized parsers for different file types.
    """
    
    file_ext = os.path.splitext(file_path)[1].lower()
    
    # Handle different file types with specialized parsers
    if file_ext in ['.xlsx', '.csv']:
        return await parse_spreadsheet_file(file_path, with_tables)
    elif file_ext == '.ipynb':
        return await parse_jupyter_notebook(file_path)
    elif file_ext in ['.ppt', '.pptx']:
        # Enable image extraction for PowerPoint files
        return await parse_powerpoint_file(file_path)
    elif file_ext == '.docx':
        # Use specialized DOCX parser
        return await parse_docx_file(file_path)
    else:
        # Use existing parsers for document files
        # Try PyMuPDF first (fast and lightweight) - enable images by default for PDFs
        text, images, tables = await parse_file_pymupdf(file_path, with_images=True, with_tables=with_tables)
        
        if text and len(text.strip()) > 50:
            # Check if extraction quality is good
            if not is_handwritten_or_poor_extraction(text):
                return text, images, tables
            else:
                pass
        else:
            pass
        
        # Fallback to Docling for better extraction
        try:
            text, images, tables = await parse_file_docling(file_path, with_images=True, with_tables=with_tables)
            if text and len(text.strip()) > 50:
                return text, images, tables
            else:
                return None, None, None
        except Exception as e:
            return None, None, None


async def process_file(
    filename: str,
) -> Union[Tuple[str, None], Tuple[None, None], Tuple[str, str]]:
    """
    Process a file locally without using LlamaCloud.
    """
    # Resolve file path - try multiple locations
    file_path = filename
    possible_paths = [
        filename,  # Try as-is
        os.path.join(os.getcwd(), filename),  # Try from current working directory
        os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), filename),  # Try from project root
    ]
    
    for path in possible_paths:
        if os.path.exists(path):
            file_path = path
            break
    else:
        return None, None
    
    text, _, _ = await parse_file(file_path=file_path)
    if text is None:
        return None, None
    
    # For now, return the text content directly
    # In the future, you could add local extraction logic here
    return "File processed successfully", text





async def query_index(question: str) -> Union[str, None]:
    """
    Query the global index (placeholder for now).
    In the future, this could query a global Pinecone index.
    """
    return f"Query: {question}\n\nThis is a placeholder response. Use query_index_for_notebook for notebook-specific queries."


async def process_file_for_notebook(
    filename: str,
    notebook_id: str,
    document_type: str = "general_review",
) -> Union[Tuple[str, None], Tuple[None, None], Tuple[str, str]]:
    """
    Process a file and create embeddings specific to a notebook using Pinecone.
    This function processes the file and adds it to the notebook's Pinecone index.
    """
    text = None
    text_chunks = None
    documents = None
    try:
        # Resolve file path - try multiple locations
        file_path = filename
        possible_paths = [
            filename,  # Try as-is
            os.path.join(os.getcwd(), filename),  # Try from current working directory
            os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), filename),  # Try from project root
        ]
        
        for path in possible_paths:
            if os.path.exists(path):
                file_path = path
                break
        else:
            print(f"File not found: {filename}")
            return None, None
        
        # Parse the file to get text content using Docling
        text, _, _ = await parse_file(file_path=file_path)
        if text is None:
            return None, None
        
        # Process text in token-aware chunks to prevent OpenAI errors
        text_chunks = smart_chunk_text(text, max_tokens=6000)
        
        # Create document dict for Pinecone with chunked content
        documents = []
        for i, chunk in enumerate(text_chunks):
            document = {
                "text": chunk,
                "filename": os.path.basename(file_path),
                "notebook_id": notebook_id,
                "chunk_index": i,
                "total_chunks": len(text_chunks),
                "processed_at": datetime.now().isoformat()
            }
            documents.append(document)
        
        # Add documents to Pinecone index for this notebook
        success = await pinecone_service.add_documents_to_notebook(
            notebook_id=notebook_id,
            documents=documents,
            metadata={"filename": os.path.basename(file_path)}
        )
        
        # Store chunk count before cleanup
        chunk_count = len(text_chunks) if text_chunks else 0
        
        # Clean up large variables immediately
        if text is not None:
            del text
        if text_chunks is not None:
            del text_chunks
        if documents is not None:
            del documents
        
        if success:
            return "Document processed and added to notebook index", f"Processed {chunk_count} chunks"
        
        return None, None
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return None, None
    finally:
        # Clean up any remaining large variables - use try/except to handle unbound variables
        try:
            if 'text' in locals() and text is not None:
                del text
        except UnboundLocalError:
            pass
        try:
            if 'text_chunks' in locals() and text_chunks is not None:
                del text_chunks
        except UnboundLocalError:
            pass
        try:
            if 'documents' in locals() and documents is not None:
                del documents
        except UnboundLocalError:
            pass
        # Force garbage collection after processing
        import gc
        gc.collect()





async def query_index_for_notebook(question: str, notebook_id: str) -> Union[str, None]:
    """
    Query the Pinecone index for a specific notebook context.
    This function queries only documents from the specified notebook.
    """
    try:
        # Query the notebook-specific Pinecone index
        raw_response = await pinecone_service.query_notebook(notebook_id, question)
        
        if not raw_response:
            return None
        
        # Format the response using template
        try:
            formatted_response = format_response_with_template(raw_response, question)
            return formatted_response
        except Exception as e:
            # Fallback to raw response
            return raw_response
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return None


async def get_plots_and_tables(
    file_path: str,
) -> Union[Tuple[Optional[List[str]], Optional[List[pd.DataFrame]]]]:
    # Resolve file path - try multiple locations
    resolved_path = file_path
    possible_paths = [
        file_path,  # Try as-is
        os.path.join(os.getcwd(), file_path),  # Try from current working directory
        os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), file_path),  # Try from project root
    ]
    
    for path in possible_paths:
        if os.path.exists(path):
            resolved_path = path
            break
    else:
        return None, None
    
    _, images, tables = await parse_file(
        file_path=resolved_path, with_images=True, with_tables=True
    )
    return images, tables


def verify_claim(
    claim: str,
    sources: str,
) -> Tuple[bool, Optional[List[str]]]:
    response = LLM_VERIFIER.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "user",
                "content": f"I have this claim: {claim} that is allegedly supported by these sources:\n\n'''\n{sources}\n'''\n\nCan you please tell me whether or not this claim is truthful and, if it is, identify one to three passages in the sources specifically supporting the claim?",
            }
        ],
        response_format={"type": "json_object"}
    )
    response_json = json.loads(response.choices[0].message.content)
    return response_json["claim_is_true"], response_json["supporting_citations"]


# Study Features Cache Functions
async def get_cached_study_feature(notebook_id: str, feature_type: str) -> Optional[str]:
    """
    Retrieve a cached study feature from the database.
    
    Args:
        notebook_id: The notebook ID
        feature_type: The type of feature ('summary', 'exam', 'flashcards')
    
    Returns:
        The cached content if found, None otherwise
    """
    try:
        result = supabase.table("study_features_cache").select("content").eq("notebook_id", notebook_id).eq("feature_type", feature_type).execute()
        
        if result.data and len(result.data) > 0:
            return result.data[0]["content"]
        return None
    except Exception as e:
        return None


async def cache_study_feature(notebook_id: str, feature_type: str, content: str) -> bool:
    """
    Cache a study feature in the database.
    
    Args:
        notebook_id: The notebook ID
        feature_type: The type of feature ('summary', 'exam', 'flashcards')
        content: The content to cache
    
    Returns:
        True if successful, False otherwise
    """
    try:
        # Use upsert to handle both insert and update cases
        result = supabase.table("study_features_cache").upsert({
            "notebook_id": notebook_id,
            "feature_type": feature_type,
            "content": content
        }).execute()
        
        return True
    except Exception as e:
        return False


async def clear_cached_study_feature(notebook_id: str, feature_type: str) -> bool:
    """
    Clear a cached study feature from the database.
    
    Args:
        notebook_id: The notebook ID
        feature_type: The type of feature ('summary', 'exam', 'flashcards')
    
    Returns:
        True if successful, False otherwise
    """
    try:
        result = supabase.table("study_features_cache").delete().eq("notebook_id", notebook_id).eq("feature_type", feature_type).execute()
        return True
    except Exception as e:
        return False


def smart_chunk_text(text: str, max_tokens: int = 6000, overlap_tokens: int = 200) -> List[str]:
    """
    Split text into chunks that respect token limits for OpenAI embeddings.
    Uses tiktoken for accurate token counting with overlapping content for context continuity.
    
    Args:
        text: Text to chunk
        max_tokens: Maximum tokens per chunk
        overlap_tokens: Number of tokens to overlap between chunks for context
    """
    try:
        import tiktoken
        
        # Use the same encoding as text-embedding-3-small
        encoding = tiktoken.encoding_for_model("text-embedding-3-small")
        
        # If text is small enough, return as single chunk
        tokens = encoding.encode(text)
        if len(tokens) <= max_tokens:
            return [text]
        
        # Split into sentences for better chunking
        sentences = re.split(r'(?<=[.!?])\s+', text)
        
        chunks = []
        current_chunk_sentences = []
        current_tokens = 0
        
        # Build chunks sentence by sentence
        for sentence in sentences:
            sentence_tokens = len(encoding.encode(sentence))
            
            # If adding this sentence would exceed the limit, finalize current chunk
            if current_tokens + sentence_tokens > max_tokens and current_chunk_sentences:
                chunk_text = " ".join(current_chunk_sentences)
                chunks.append(chunk_text)
                
                # Create overlap for next chunk
                overlap_sentences = []
                overlap_token_count = 0
                
                # Take sentences from the end of current chunk for overlap
                for i in range(len(current_chunk_sentences) - 1, -1, -1):
                    overlap_sentence = current_chunk_sentences[i]
                    overlap_sentence_tokens = len(encoding.encode(overlap_sentence))
                    
                    if overlap_token_count + overlap_sentence_tokens <= overlap_tokens:
                        overlap_sentences.insert(0, overlap_sentence)
                        overlap_token_count += overlap_sentence_tokens
                    else:
                        break
                
                # Start new chunk with overlap + current sentence
                current_chunk_sentences = overlap_sentences + [sentence]
                current_tokens = overlap_token_count + sentence_tokens
            else:
                current_chunk_sentences.append(sentence)
                current_tokens += sentence_tokens
        
        # Add the last chunk if it has content
        if current_chunk_sentences:
            chunk_text = " ".join(current_chunk_sentences)
            chunks.append(chunk_text)
        
        # Handle edge case where a single sentence is too long
        final_chunks = []
        for chunk in chunks:
            chunk_tokens = len(encoding.encode(chunk))
            if chunk_tokens <= max_tokens:
                final_chunks.append(chunk)
            else:
                # Split by words if sentence is too long
                words = chunk.split()
                temp_chunk_words = []
                temp_tokens = 0
                
                for word in words:
                    word_tokens = len(encoding.encode(word + " "))
                    if temp_tokens + word_tokens > max_tokens and temp_chunk_words:
                        # Create word-based chunk
                        word_chunk = " ".join(temp_chunk_words)
                        final_chunks.append(word_chunk)
                        
                        # Create overlap for next word chunk
                        overlap_words = temp_chunk_words[-50:] if len(temp_chunk_words) > 50 else temp_chunk_words[-len(temp_chunk_words)//2:]
                        overlap_word_tokens = len(encoding.encode(" ".join(overlap_words)))
                        
                        if overlap_word_tokens <= overlap_tokens:
                            temp_chunk_words = overlap_words + [word]
                            temp_tokens = overlap_word_tokens + word_tokens
                        else:
                            temp_chunk_words = [word]
                            temp_tokens = word_tokens
                    else:
                        temp_chunk_words.append(word)
                        temp_tokens += word_tokens
                
                if temp_chunk_words:
                    final_chunks.append(" ".join(temp_chunk_words))
        
        return final_chunks if final_chunks else [text[:max_tokens*3]]
        
    except ImportError:
        # Fallback to character-based chunking with overlap if tiktoken is not available
        chunk_size = max_tokens * 3  # Rough estimate: 1 token ≈ 3-4 characters
        overlap_size = overlap_tokens * 3
        chunks = []
        
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            chunks.append(chunk)
            
            # Move start forward but with overlap
            start = end - overlap_size if end - overlap_size > start else end
            
            if start >= len(text):
                break
                
        return chunks
    except Exception as e:
        # Fallback to simple character chunking with overlap on any error
        chunk_size = max_tokens * 3
        overlap_size = overlap_tokens * 3
        chunks = []
        
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            chunks.append(chunk)
            
            # Move start forward but with overlap
            start = end - overlap_size if end - overlap_size > start else end
            
            if start >= len(text):
                break
                
        return chunks
